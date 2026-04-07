from __future__ import annotations

import argparse
import re
import subprocess
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def count_slides(html: str) -> int:
    return len(re.findall(r'<section[^>]*class="[^"]*\bslide\b', html))


def chrome_path() -> str:
    candidates = [
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    raise FileNotFoundError("Chrome/Chromium not found in /Applications")


def run() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--width", type=int, default=1920)
    parser.add_argument("--height", type=int, default=1080)
    parser.add_argument("--capture-height", type=int, default=1165)
    parser.add_argument("--time-budget", type=int, default=9000)
    parser.add_argument("--timeout", type=int, default=45)
    args = parser.parse_args()

    root = Path(__file__).resolve().parent
    html_path = root / "index.html"
    html = html_path.read_text(encoding="utf-8", errors="ignore")
    n = count_slides(html)
    if n <= 0:
        raise RuntimeError("No slides found")

    out_dir = root / ".export"
    out_dir.mkdir(exist_ok=True)

    chrome = chrome_path()
    base_url = f"file://{html_path.as_posix()}"

    raw_paths: list[Path] = []
    slide_paths: list[Path] = []
    for i in range(1, n + 1):
        print(f"Render {i}/{n}")
        raw = out_dir / f"raw-{i:02d}.png"
        url = f"{base_url}?export=1&slide={i}"
        subprocess.run(
            [
                chrome,
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                f"--window-size={args.width},{args.capture_height}",
                f"--virtual-time-budget={args.time_budget}",
                f"--screenshot={raw.as_posix()}",
                url,
            ],
            check=True,
            timeout=args.timeout,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        raw_paths.append(raw)

        img = Image.open(raw).convert("RGB").crop((0, 0, args.width, args.height))
        slide = out_dir / f"slide-{i:02d}.png"
        img.save(slide)
        slide_paths.append(slide)

    imgs = [Image.open(p).convert("RGB") for p in slide_paths]
    pdf_path = root / "IXI.pdf"
    imgs[0].save(pdf_path.as_posix(), save_all=True, append_images=imgs[1:])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in slide_paths:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(p.as_posix(), 0, 0, width=prs.slide_width, height=prs.slide_height)
    pptx_path = root / "IXI.pptx"
    prs.save(pptx_path.as_posix())

    return 0


if __name__ == "__main__":
    raise SystemExit(run())
